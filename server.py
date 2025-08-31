#server.py
import os
import io
import json
import csv
import re
import logging
import traceback
from datetime import datetime
from typing import Tuple, List
from functools import wraps

import requests
from flask import Flask, jsonify, send_from_directory, request, session, redirect, url_for
from flask_cors import CORS
from dotenv import load_dotenv
import chromadb
from chromadb.config import Settings
from openai import OpenAI

# --- Опциональные парсеры (мягкие зависимости) ---
try:
    import pypdf  # для PDF
except Exception:
    pypdf = None

try:
    from docx import Document  # для DOCX
except Exception:
    Document = None

try:
    from striprtf.striprtf import rtf_to_text  # для RTF
except Exception:
    rtf_to_text = None

try:
    from pptx import Presentation  # для PPTX
except Exception:
    Presentation = None

# --- Конфигурация ---
load_dotenv()

# Настройка логирования
log_level = getattr(logging, os.getenv('LOG_LEVEL', 'INFO').upper(), logging.INFO)
log_file = os.getenv('LOG_FILE', 'app.log')

logging.basicConfig(
    level=log_level if os.getenv('DEBUG') != 'true' else logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


# Дефолтные значения
DEFAULT_CHUNK_SIZE = 1000
DEFAULT_CHUNK_OVERLAP = 200
DEFAULT_MAX_FILE_BYTES = 15

# Получаем значения из .env
CHUNK_SIZE = int(os.getenv('CHUNK_SIZE', DEFAULT_CHUNK_SIZE))
CHUNK_OVERLAP = int(os.getenv('CHUNK_OVERLAP', DEFAULT_CHUNK_OVERLAP))
MAX_FILE_BYTES = int(os.getenv('MAX_UPLOAD_SIZE_MB', DEFAULT_MAX_FILE_BYTES)) * 1024 * 1024

# Получаем разрешенные расширения файлов из .env
ALLOWED_EXT_STRING = os.getenv('ALLOWED_FILE_EXTENSIONS', 'txt,md,pdf,docx,rtf,csv,json,pptx')
ALLOWED_EXT = {f'.{ext.strip()}' for ext in ALLOWED_EXT_STRING.split(',') if ext.strip()}

# Конфигурация для продакшена
IS_PRODUCTION = os.getenv('FLASK_ENV') == 'production'
SECRET_KEY = os.getenv('SECRET_KEY', 'dev-secret-key-change-in-production')
if IS_PRODUCTION and SECRET_KEY == 'dev-secret-key-change-in-production':
    logger.error("ВНИМАНИЕ: Используется дефолтный SECRET_KEY в продакшене!")
    
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
if not OPENAI_API_KEY:
    logger.error("OPENAI_API_KEY не найден в переменных окружения")
    raise SystemExit("Необходимо задать OPENAI_API_KEY")

# --- Инициализация клиентов ---
try:
    openai_client = OpenAI(api_key=OPENAI_API_KEY)
    logger.info("OpenAI клиент успешно инициализирован")
except Exception as e:
    logger.error(f"Не удалось инициализировать OpenAI клиент: {e}")
    raise SystemExit(1)

try:
    chroma_db_path = os.getenv('CHROMA_DB_PATH', 'chroma_db')
    chroma_telemetry = os.getenv('CHROMA_ANONYMIZED_TELEMETRY', 'false').lower() == 'true'
    chroma_client = chromadb.PersistentClient(
        path=chroma_db_path,
        settings=Settings(anonymized_telemetry=chroma_telemetry)
    )
except Exception as e:
    logger.error(f"Не удалось инициализировать ChromaDB клиент: {e}")
    raise SystemExit(1)

# --- Функции для работы с пользовательскими коллекциями ---
def get_user_collection_name(user_id):
    return f"knowledge_base_user_{user_id}"

def get_user_collection(user_id):
    collection_name = get_user_collection_name(user_id)
    return chroma_client.get_or_create_collection(name=collection_name)

# --- Хелперы ---
def split_text_into_chunks(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
    if not text:
        return []
    chunks = []
    start = 0
    step = max(1, chunk_size - chunk_overlap)
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += step
    return chunks

def norm_ws(s: str) -> str:
    return re.sub(r'[ \t\f\v]+', ' ', re.sub(r'\r\n?', '\n', s)).strip()

def file_ext(name: str) -> str:
    return os.path.splitext(name or '')[1].lower()

def decode_text_bytes(b: bytes) -> str:
    # предпочтительно utf-8, но терпимо к ошибкам
    try:
        return b.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return b.decode('utf-8', errors='ignore')
        except Exception:
            return b.decode('latin-1', errors='ignore')

def extract_from_pdf(data: bytes) -> Tuple[str, str]:
    if not pypdf:
        return "", "нужен пакет 'pypdf' для PDF"
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        texts = []
        for page in reader.pages:
            t = page.extract_text() or ""
            if t:
                texts.append(t)
        return norm_ws("\n\n".join(texts)), ""
    except Exception as e:
        return "", f'ошибка PDF: {e}'

def extract_from_docx(data: bytes) -> Tuple[str, str]:
    if not Document:
        return "", "нужен пакет 'python-docx' для DOCX"
    try:
        doc = Document(io.BytesIO(data))
        parts = []
        for p in doc.paragraphs:
            if p.text:
                parts.append(p.text)
        # Таблицы
        for tbl in getattr(doc, 'tables', []):
            for row in tbl.rows:
                row_text = " | ".join(cell.text or "" for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return norm_ws("\n".join(parts)), ""
    except Exception as e:
        return "", f'ошибка DOCX: {e}'

def extract_from_rtf(data: bytes) -> Tuple[str, str]:
    if rtf_to_text:
        try:
            return norm_ws(rtf_to_text(decode_text_bytes(data))), ""
        except Exception as e:
            return "", f'ошибка RTF: {e}'
    # наивное удаление управляющих последовательностей
    try:
        s = decode_text_bytes(data)
        s = re.sub(r'\\[a-z]+\d* ?|{\\[^}]+}|[{}]', ' ', s, flags=re.I)
        return norm_ws(s), "упрощённый парсер RTF (лучше установить striprtf)"
    except Exception as e:
        return "", f'ошибка RTF: {e}'

def extract_from_csv(data: bytes) -> Tuple[str, str]:
    try:
        s = decode_text_bytes(data)
        out_lines = []
        reader = csv.reader(io.StringIO(s))
        for row in reader:
            out_lines.append(" | ".join(row))
        return norm_ws("\n".join(out_lines)), ""
    except Exception as e:
        return "", f'ошибка CSV: {e}'

def flatten_json(obj, prefix="") -> List[str]:
    lines = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            new_prefix = f"{prefix}.{k}" if prefix else str(k)
            lines.extend(flatten_json(v, new_prefix))
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            new_prefix = f"{prefix}[{i}]"
            lines.extend(flatten_json(v, new_prefix))
    else:
        lines.append(f"{prefix}: {obj}")
    return lines

def extract_from_json(data: bytes) -> Tuple[str, str]:
    try:
        obj = json.loads(decode_text_bytes(data))
        lines = flatten_json(obj)
        return norm_ws("\n".join(lines)), ""
    except Exception as e:
        return "", f'ошибка JSON: {e}'

def extract_from_md_or_txt(data: bytes) -> Tuple[str, str]:
    s = decode_text_bytes(data)
    # упрощённое «снятие» части markdown-разметки
    s = re.sub(r'`{1,3}[\s\S]*?`{1,3}', ' ', s)  # кодовые блоки
    s = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', s)  # ссылки [текст](url)
    s = re.sub(r'#{1,6}\s*', '', s)  # заголовки
    return norm_ws(s), ""

def extract_from_pptx(data: bytes) -> Tuple[str, str]:
    """
    Извлекаем текст из .pptx: текстовые блоки, таблицы и (по возможности) заметки к слайдам.
    Требуется пакет python-pptx.
    """
    if not Presentation:
        return "", "нужен пакет 'python-pptx' для PPTX"
    try:
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            # Текстовые шейпы
            for shape in slide.shapes:
                # Текст в shape
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = getattr(shape, "text", "") or ""
                    if text.strip():
                        parts.append(text)
                # Таблицы
                if hasattr(shape, "has_table") and shape.has_table:
                    try:
                        tbl = shape.table
                        for row in tbl.rows:
                            row_text = " | ".join(cell.text or "" for cell in row.cells)
                            if row_text.strip():
                                parts.append(row_text)
                    except Exception:
                        pass
            # Заметки к слайдам (если есть)
            try:
                if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                    notes_tf = getattr(slide.notes_slide, "notes_text_frame", None)
                    if notes_tf:
                        note_text = "\n".join(p.text or "" for p in notes_tf.paragraphs)
                        if note_text.strip():
                            parts.append(f"[Notes] {note_text}")
            except Exception:
                pass
        return norm_ws("\n".join(parts)), ""
    except Exception as e:
        return "", f'ошибка PPTX: {e}'

def read_file_to_text(filename: str, data: bytes) -> Tuple[str, str]:
    ext = file_ext(filename)
    if ext not in ALLOWED_EXT:
        return "", "неподдерживаемый формат"
    if len(data) > MAX_FILE_BYTES:
        return "", "слишком большой файл"
    if ext in ('.txt', '.md'):
        return extract_from_md_or_txt(data)
    if ext == '.pdf':
        return extract_from_pdf(data)
    if ext == '.docx':
        return extract_from_docx(data)
    if ext == '.rtf':
        return extract_from_rtf(data)
    if ext == '.csv':
        return extract_from_csv(data)
    if ext == '.json':
        return extract_from_json(data)
    if ext == '.pptx':
        return extract_from_pptx(data)
    return "", "неизвестный формат"

# --- Flask ---
app = Flask(__name__, static_folder='public', static_url_path='')
app.secret_key = SECRET_KEY
# CORS настройка из переменной окружения
allowed_origins = os.getenv('ALLOWED_ORIGINS', 'http://localhost:3000,https://your-domain.com').split(',') if IS_PRODUCTION else ['*']
CORS(app, origins=[origin.strip() for origin in allowed_origins])

# Декоратор для логирования ошибок
def log_errors(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        try:
            return f(*args, **kwargs)
        except Exception as e:
            logger.error(f"Ошибка в {f.__name__}: {str(e)}")
            logger.error(f"Traceback: {traceback.format_exc()}")
            return jsonify({'error': 'Внутренняя ошибка сервера'}), 500
    return decorated_function

# Декоратор для логирования API запросов
def log_requests(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        user_id = session.get('user_id', 'anonymous')
        logger.info(f"API запрос: {request.method} {request.path} от пользователя {user_id}")
        start_time = datetime.now()
        result = f(*args, **kwargs)
        end_time = datetime.now()
        duration = (end_time - start_time).total_seconds()
        logger.info(f"API ответ: {request.path} выполнен за {duration:.2f}с")
        return result
    return decorated_function

# --- Функции для работы с пользователями ---
def load_users():
    users_file = os.getenv('USERS_FILE', 'users.json')
    try:
        with open(users_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        return []

def authenticate_user(login, password):
    users = load_users()
    for user in users:
        if user['login'] == login and user['password'] == password:
            return {'id': user['id'], 'login': user['login']}
    return None

def require_auth():
    if 'user_id' not in session:
        return False
    return True

# --- Роуты ---
@app.route('/login', methods=['GET', 'POST'])
@log_requests
@log_errors
def login():
    if request.method == 'GET':
        return send_from_directory('public', 'login.html')
    
    data = request.get_json()
    if not data or 'login' not in data or 'password' not in data:
        return jsonify({'error': 'Не указаны логин или пароль'}), 400
    
    user = authenticate_user(data['login'], data['password'])
    if user:
        session['user_id'] = user['id']
        session['user_login'] = user['login']
        logger.info(f"Успешная авторизация пользователя: {user['login']} (ID: {user['id']})")
        return jsonify({'success': True, 'user': user}), 200
    else:
        logger.warning(f"Неудачная попытка авторизации для логина: {data['login']}")
        return jsonify({'error': 'Неверный логин или пароль'}), 401

@app.route('/logout', methods=['POST'])
def logout():
    session.clear()
    return jsonify({'success': True}), 200

@app.route('/')
def index():
    if not require_auth():
        return redirect(url_for('login'))
    return send_from_directory('public', 'agent.html')

@app.route('/upload', methods=['POST'])
@log_requests
@log_errors
def upload_files():
    if not require_auth():
        return jsonify({'error': 'Требуется авторизация'}), 401
    
    user_id = session.get('user_id')
    collection = get_user_collection(user_id)
    
    if 'files' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400
    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({'error': 'No selected files'}), 400

    try:
        all_chunks, all_metadatas, all_ids = [], [], []
        warnings = []
        indexed_files = 0
        skipped_files = 0

        for file in files:
            filename = file.filename
            ext = file_ext(filename)
            if ext not in ALLOWED_EXT:
                warnings.append(f"{filename}: формат не поддерживается")
                skipped_files += 1
                continue

            blob = file.read()
            text, note = read_file_to_text(filename, blob)

            if note:
                warnings.append(f"{filename}: {note}")

            if not text or not text.strip():
                warnings.append(f"{filename}: не удалось извлечь печатный текст")
                skipped_files += 1
                continue

            # Чанкуем
            chunks = split_text_into_chunks(text, CHUNK_SIZE, CHUNK_OVERLAP)
            if not chunks:
                warnings.append(f"{filename}: пустой документ")
                skipped_files += 1
                continue

            for i, chunk in enumerate(chunks):
                uid = f"{filename}-{i}"
                all_chunks.append(chunk)
                all_metadatas.append({'filename': filename, 'chunk_number': i})
                all_ids.append(uid)
            indexed_files += 1

        if not all_chunks:
            msg = 'Нет данных для индексации.'
            if warnings:
                msg += ' ' + '; '.join(warnings)
            return jsonify({'message': msg, 'warnings': warnings}), 200

        # Эмбеддинги
        embedding_model = os.getenv('OPENAI_MODEL_EMBEDDING', 'text-embedding-3-small')
        embedding_response = openai_client.embeddings.create(
            input=all_chunks, model=embedding_model
        )
        embeddings = [e.embedding for e in embedding_response.data]
        collection.add(ids=all_ids, embeddings=embeddings, documents=all_chunks, metadatas=all_metadatas)

        total_chunks_indexed = len(all_chunks)
        msg = f'Успешно проиндексировано {total_chunks_indexed} фрагментов из {indexed_files} файлов.'
        if skipped_files:
            msg += f' Пропущено файлов: {skipped_files}.'
        return jsonify({'message': msg, 'warnings': warnings}), 200

    except Exception as e:
        print(f"Error during file indexing: {e}")
        return jsonify({'error': f'Произошла ошибка при индексации: {str(e)}'}), 500

@app.route('/query', methods=['POST'])
def query_knowledge_base():
    if not require_auth():
        return jsonify({'error': 'Требуется авторизация'}), 401
    
    user_id = session.get('user_id')
    collection = get_user_collection(user_id)
    
    try:
        data = request.get_json(force=True) or {}
        query_text = (data.get('query') or '').strip()
        if not query_text:
            return jsonify({'error': 'Query text is missing'}), 400
        try:
            default_n_results = int(os.getenv('DEFAULT_N_RESULTS', 3))
            n_results = int(data.get('n_results', default_n_results))
        except Exception:
            n_results = int(os.getenv('DEFAULT_N_RESULTS', 3))
        
        min_n_results = int(os.getenv('MIN_N_RESULTS', 1))
        max_n_results = int(os.getenv('MAX_N_RESULTS', 10))
        n_results = max(min_n_results, min(max_n_results, n_results))

        embedding_model = os.getenv('OPENAI_MODEL_EMBEDDING', 'text-embedding-3-small')
        query_embedding = openai_client.embeddings.create(
            input=[query_text], model=embedding_model
        ).data[0].embedding

        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=n_results,
            include=["documents", "metadatas", "distances"]
        )

        response_data = []
        if results and results.get('ids') and results['ids'][0]:
            for i in range(len(results['ids'][0])):
                response_data.append({
                    'id': results['ids'][0][i],
                    'document': results['documents'][0][i],
                    'metadata': results['metadatas'][0][i],
                    'distance': results['distances'][0][i]
                })
        return jsonify(response_data)
    except Exception as e:
        print(f"Error during query: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/files', methods=['GET'])
@log_requests
@log_errors
def get_indexed_files():
    if not require_auth():
        return jsonify({'error': 'Требуется авторизация'}), 401
    
    user_id = session.get('user_id')
    collection = get_user_collection(user_id)
    
    try:
        results = collection.get(include=['metadatas'])
        if not results or 'metadatas' not in results or not results['metadatas']:
            return jsonify([])
        
        filenames = {meta.get('filename') for meta in results['metadatas'] if meta and meta.get('filename')}
        filenames.discard(None)  # Удаляем None значения
        return jsonify(sorted(list(filenames)))
    except Exception as e:
        logger.error(f"Ошибка получения списка файлов для пользователя {user_id}: {e}")
        return jsonify({'error': 'Не удалось загрузить список файлов'}), 500

@app.route('/reset-knowledge-base', methods=['DELETE'])
def reset_knowledge_base():
    if not require_auth():
        return jsonify({'error': 'Требуется авторизация'}), 401
    
    user_id = session.get('user_id')
    collection_name = get_user_collection_name(user_id)
    
    try:
        chroma_client.delete_collection(name=collection_name)
        # Коллекция будет автоматически создана при следующем обращении
        return jsonify({'message': 'База знаний успешно очищена.'}), 200
    except Exception as e:
        print(f"Error resetting knowledge base: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/session', methods=['GET'])
@log_requests
@log_errors
def create_session():
    if not require_auth():
        return jsonify({'error': 'Требуется авторизация'}), 401
    try:
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            return jsonify({'error': 'OPENAI_API_KEY not found in .env file'}), 500
        headers = {'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'}
        
        # Получаем конфигурацию из переменных окружения
        realtime_model = os.getenv('OPENAI_MODEL_REALTIME', 'gpt-realtime')
        voice = os.getenv('OPENAI_VOICE', 'marin')
        modalities = os.getenv('OPENAI_MODALITIES', 'audio,text').split(',') 
        instructions = os.getenv('OPENAI_INSTRUCTIONS', 'Отвечай по-русски, дружелюбно и кратко. Если пользователя не слышно — вежливо попроси повторить.')
        
        payload = {
            'model': realtime_model,
            'voice': voice,
            'modalities': [m.strip() for m in modalities],
            'instructions': instructions
        }
        api_base_url = os.getenv('OPENAI_API_BASE_URL', 'https://api.openai.com/v1')
        timeout = int(os.getenv('REALTIME_SESSION_TIMEOUT', 20))
        response = requests.post(f'{api_base_url}/realtime/sessions', headers=headers, json=payload, timeout=timeout)
        response.raise_for_status()
        data = response.json()
        client_secret = data.get('client_secret', {})
        secret_value = client_secret.get('value')
        expires_at = client_secret.get('expires_at')
        if not secret_value:
            return jsonify({'error': 'Failed to get client_secret from OpenAI response', 'details': data}), 500
        return jsonify({'client_secret': secret_value, 'expires_at': expires_at})
    except requests.exceptions.HTTPError as http_err:
        return jsonify({'error': 'Failed to create realtime session', 'details': str(http_err),
                        'response_text': response.text}), 500
    except Exception as err:
        return jsonify({'error': str(err)}), 500

if __name__ == '__main__':

    port = int(os.getenv('PORT', 3000))
    host = os.getenv('HOST', '0.0.0.0')
    debug = os.getenv('DEBUG', 'true').lower() == 'true'
    app.run(host=host, port=port, debug=debug)